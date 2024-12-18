import base64
import random
from datetime import datetime
import os
import sqlite3
import requests
import io
import numpy as np
from flask import Flask, render_template, request, jsonify, send_file, session
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
from bs4 import BeautifulSoup
from ollama import chat
from ollama import ChatResponse
import subprocess
import smtplib
from email.mime.text import MIMEText
from deep_translator import GoogleTranslator
from flask import render_template, request, redirect, url_for, flash
import bcrypt
2
app = Flask(__name__, static_folder="static", template_folder="templates")
app.secret_key = "yehbbcf83725gfbefuwfew08348gsfe732"
current_datetime = datetime.now()


def hash_password(password):
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')


def check_password(hashed, plain):
    return bcrypt.checkpw(plain.encode('utf-8'), hashed.encode('utf-8'))


def initialize_db():
    conn = sqlite3.connect("db/users.db")
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_name TEXT UNIQUE NOT NULL,
            hash_password TEXT NOT NULL,
            history_of_request TEXT DEFAULT ''
        )
    """)
    conn.commit()
    conn.close()


def gpt_response(subject, century):
    response: ChatResponse = chat(model='llama3.2', messages=[
        {
            'role': 'user',
            'content': f"Tell us about {subject} in {century}century (how it was created, what materials it was made of), in 30 words",
        },
    ])

    return response.message.content


def google_search(query):
    url = f"https://www.google.com/search?q={query}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, "html.parser")
        first_result = soup.select_one(".tF2Cxc")
        if first_result:
            title = first_result.select_one(".DKV0Md").text
            link = first_result.select_one(".yuRUbf a")["href"]
            snippet = first_result.select_one(".VwiC3b").text
            return {"title": title, "link": link, "snippet": snippet}
        else:
            return "Результатов не найдено"
    else:
        return f"Ошибка при выполнении запроса: {response.status_code}"


def remove_bg(image_path, threshold=240):
    image = Image.open(image_path).convert("RGBA")
    data = np.array(image)
    mask = (data[:, :, :3].sum(axis=2) > threshold * 3)
    data[mask, 3] = 0
    result_image = Image.fromarray(data, mode="RGBA")
    filename, file_extension = os.path.splitext(image_path)
    result_image.save(filename + '.png')


def add_resized_picture(slide, img_path, max_width, max_height, left, top):
    img = Image.open(img_path)
    width, height = img.size

    scale = min(max_width / width, max_height / height)
    new_width = int(width * scale)
    new_height = int(height * scale)

    slide.shapes.add_picture(img_path, left, top, width=Inches(new_width / 96), height=Inches(new_height / 96))


def search_exact_items_as_dict(item_name):
    conn = sqlite3.connect('db/items.db')
    cursor = conn.cursor()

    cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
    rows = cursor.fetchall()

    items = []
    for row in rows:
        item_dict = {
            "century": row[0],
            "item_name": row[1],
            "image": row[2]
        }
        items.append(item_dict)

    conn.close()
    if not items:
        conn = sqlite3.connect('db/images.db')
        cursor = conn.cursor()

        cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
        rows = cursor.fetchall()

        items = []
        for row in rows:
            item_dict = {
                "century": row[0],
                "item_name": row[1],
                "image": row[2]
            }
            items.append(item_dict)
    if not items:
        conn = sqlite3.connect('db/images.db')
        cursor = conn.cursor()

        cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
        rows = cursor.fetchall()

        items = []
        for row in rows:
            item_dict = {
                "century": row[0],
                "item_name": row[1],
                "image": row[2]
            }
            items.append(item_dict)
    if not items:
        conn = sqlite3.connect('db/images2.db')
        cursor = conn.cursor()

        cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
        rows = cursor.fetchall()

        items = []
        for row in rows:
            item_dict = {
                "century": row[0],
                "item_name": row[1],
                "image": row[2]
            }
            items.append(item_dict)
    if not items:
        conn = sqlite3.connect('db/items_database (10).db')
        cursor = conn.cursor()

        cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
        rows = cursor.fetchall()

        items = []
        for row in rows:
            item_dict = {
                "century": row[0],
                "item_name": row[1],
                "image": row[2]
            }
            items.append(item_dict)
    if not items:
        conn = sqlite3.connect('db/items_database (11).db')
        cursor = conn.cursor()

        cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
        rows = cursor.fetchall()

        items = []
        for row in rows:
            item_dict = {
                "century": row[0],
                "item_name": row[1],
                "image": row[2]
            }
            items.append(item_dict)
    return items


def int_to_roman(n):
    roman_numerals = {
        1: "I", 4: "IV", 5: "V", 9: "IX", 10: "X",
        40: "XL", 50: "L", 90: "XC", 100: "C",
        400: "CD", 500: "D", 900: "CM", 1000: "M"
    }
    result = ""
    for value in sorted(roman_numerals.keys(), reverse=True):
        while n >= value:
            result += roman_numerals[value]
            n -= value
    return result


def save_user_to_file(username, password):
    with open("users.txt", "a") as file:
        file.write(f"{username},{password}\n")


def get_users_from_file():
    if not os.path.exists("users.txt"):
        return []
    with open("users.txt", "r") as file:
        lines = file.readlines()
    return [line.strip().split(",") for line in lines]


def authenticate_user(user_name, password):
    conn = sqlite3.connect("db/users.db")
    cursor = conn.cursor()
    cursor.execute("SELECT hash_password FROM users WHERE user_name = ?", (user_name,))
    row = cursor.fetchone()
    conn.close()
    if row and check_password(row[0], password):
        return True
    return False


def update_history(user_name, request):
    conn = sqlite3.connect("db/users.db")
    cursor = conn.cursor()
    formatted_date = current_datetime.strftime("%d.%m.%Y")
    hours_minutes_seconds = current_datetime.strftime("%H:%M:%S")
    cursor.execute("""
        UPDATE users
        SET history_of_request = history_of_request || ?
        WHERE user_name = ?
    """, (f"{request},{formatted_date} {hours_minutes_seconds};", user_name))
    conn.commit()
    conn.close()


def register_user(user_name, password, email):
    conn = sqlite3.connect("db/users.db")
    cursor = conn.cursor()
    hash_pass = hash_password(password)
    try:
        cursor.execute("INSERT INTO users (user_name, hash_password, email) VALUES (?, ?, ?)",
                       (user_name, hash_pass, email))
        conn.commit()
    except sqlite3.IntegrityError:
        return False  # Username already exists
    finally:
        conn.close()
    return True


def change_password(email, password):
    conn = sqlite3.connect("db/users.db")
    cursor = conn.cursor()
    hash_pass = hash_password(password)
    cursor.execute("""
                    UPDATE users
                    SET hash_password = ?
                    WHERE email = ?
                """, (hash_pass, email))
    conn.commit()


def create_of_presentation(item_name):
    a = item_name
    prs = Presentation()
    search_query = a
    items = sorted(search_exact_items_as_dict(search_query), key=lambda x: int(x["century"]))

    c = 0
    for i in items:
        c += 1
        print(i)
        century = i["century"]
        sub = i["item_name"]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        img_path = 'img1/back.png'
        slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)
        add_resized_picture(slide, "img1/img_1.png", 1 * 96, 1 * 96, left=Inches(9), top=Inches(6.5))
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title = title_box.text_frame
        title.text = f"{sub} {century} век"
        title.paragraphs[0].font.size = Pt(36)
        title.paragraphs[0].font.bold = True
        title.paragraphs[0].font.name = "Times New Roman"
        title.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4), Inches(8))
        text = text_box.text_frame
        f = gpt_response(sub, century)
        if session.get('language') == 'ru':
            f = GoogleTranslator(source='en', target='ru').translate(f)
        text.text = f
        text.paragraphs[0].font.size = Pt(18)
        text.paragraphs[0].font.name = "Times New Roman"
        text.word_wrap = True
        image = Image.open(io.BytesIO(i["image"]))
        image_path = f"{sub}.png"
        image.save(image_path)
        max_width = 6 * 96
        max_height = 3 * 96
        add_resized_picture(slide, "img1/icon.png", 1 * 56, 1 * 56, left=Inches(4.5), top=Inches(1.95))
        add_resized_picture(slide, image_path, max_width, max_height, left=Inches(1.2), top=Inches(2))
        os.remove(image_path)

        start_x = Inches(1)
        y_position = Inches(5.5)
        line_length = Inches(8)
        num_centuries = 5
        century_spacing = line_length / (num_centuries - 1)
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, y_position, start_x + line_length, y_position
        )
        line_shape.line.width = Pt(2)
        line_shape.line.color.rgb = RGBColor(0, 0, 0)

        centuries = [
            (f"{int_to_roman(int(century) - 2)} век"),
            (f"{int_to_roman(int(century) - 1)} век"),
            (f"{int_to_roman(int(century))} век"),
            (f"{int_to_roman(int(century) + 1)} век"),
            (f"{int_to_roman(int(century) + 2)} век")
        ]

        for i, century in enumerate(centuries):
            x_position = start_x + i * century_spacing

            century_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left=x_position,
                top=y_position - Inches(0.1),
                width=Pt(12),
                height=Pt(12)
            )
            century_shape.fill.solid()
            if i == 0 or i == 1 or i == 3 or i == 4:
                century_shape.fill.fore_color.rgb = RGBColor(0, 0, 255)
            else:
                century_shape.fill.fore_color.rgb = RGBColor(220, 20, 60)
            textbox = slide.shapes.add_textbox(
                left=x_position - Inches(0.3),
                top=y_position + Inches(0.3),
                width=Inches(1),
                height=Inches(0.5)
            )
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = century
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    g = '1.pptx'
    prs.save(g)
    return g


def get_user_history(username):
    conn = sqlite3.connect('db/users.db')
    cursor = conn.cursor()
    cursor.execute("SELECT history_of_request FROM users WHERE user_name = ?", (username,))
    result = cursor.fetchone()

    conn.close()
    return result[0]


def clear_user_history(username):
    conn = sqlite3.connect('db/users.db')
    cursor = conn.cursor()

    # Устанавливаем значение поля history_of_request пустым
    cursor.execute("UPDATE users SET history_of_request = '' WHERE user_name = ?", (username,))
    conn.commit()
    conn.close()


@app.route('/templates/create_of_presentation.html')
def create_of_pr():
    item_name = request.args.get('item_name', '')
    if not item_name:
        return "Item name not provided.", 400
    create_of_presentation(item_name)
    file_path = '1.pptx'
    return send_file(file_path, as_attachment=True, download_name="presentation.pptx")


@app.route('/')
def main_page():
    username = session.get('username')
    return render_template("main-page.html", username=username)


@app.route('/update_favorites', methods=['POST'])
def update_favorites():
    data = request.json
    username = session.get('username')  # Получить текущего пользователя из сессии
    item_name = data.get('item_name')  # Имя предмета

    if not username or not item_name:
        return jsonify({'status': 'error', 'message': 'Username or item name is missing'}), 400

    try:
        conn = sqlite3.connect('db/users.db')
        cursor = conn.cursor()

        # Получить текущий список избранного
        cursor.execute("SELECT favorites FROM users WHERE user_name = ?", (username,))
        row = cursor.fetchone()
        if row:
            favorites = row[0].split(',') if row[0] else []

            # Добавить или удалить предмет
            if item_name in favorites:
                favorites.remove(item_name)
            else:
                favorites.append(item_name)

            # Обновить запись в базе данных
            updated_favorites = ','.join(favorites)
            cursor.execute("UPDATE users SET favorites = ? WHERE user_name = ?", (updated_favorites, username))
            conn.commit()

        conn.close()
        return jsonify({'status': 'success', 'favorites': favorites})
    except Exception as e:
        print(f"Error updating favorites: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/templates/history.html/clear_history', methods=['POST'])
def clear_history():
    username = session.get('username')
    if not username:
        flash("Вы должны быть авторизованы, чтобы очищать историю.")
        return redirect('/login')
    clear_user_history(username)
    return render_template('history.html', username=username)


@app.route('/templates/main-page.html')
def main_page1():
    username = session.get('username')
    return render_template("main-page.html", username=username)


@app.route('/templates/authors.html')
def authors_page():
    username = session.get('username')
    return render_template("authors.html", username=username)


@app.route('/templates/history.html')
def history_page():
    username = session.get('username')
    user_history = get_user_history(username).split(";")
    user_history = [i.split(",") for i in user_history if i != ""]
    return render_template("history.html", username=username, user_history=user_history)


@app.route('/templates/review.html')
def review_page():
    username = session.get('username')
    item_name = request.args.get('input_review', '')
    print(item_name)
    return render_template("review.html", username=username)


def send_email(smtp_server, smtp_port, sender_email, sender_password, recipient_email, subject, body):
    """
    Простая функция для отправки письма через SMTP с поддержкой не-ASCII символов.
    """
    # Формирование сообщения с правильной кодировкой
    message = MIMEText(body, "plain", "utf-8")  # Указываем кодировку utf-8
    message["Subject"] = subject
    message["From"] = sender_email
    message["To"] = recipient_email

    try:
        # Установка соединения с сервером
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()  # Включение защищенного соединения
            server.login(sender_email, sender_password)  # Логин
            server.sendmail(sender_email, recipient_email, message.as_string())  # Отправка письма
        print("Письмо успешно отправлено!")
    except Exception as e:
        print(f"Ошибка при отправке письма: {e}")


@app.route('/submit-review', methods=['POST'])
def submit_review():
    username = session.get('username')
    review_text = request.form.get('review_text')
    smtp_server = "smtp.gmail.com"
    smtp_port = 587
    sender_email = "dictionary.of.evolutions@gmail.com"
    app_password = "jzje ysgq plxm oulg"
    recipient_email = "hugofinger055@gmail.com"
    subject = "Отзыв о нашем сайте"
    body = review_text
    send_email(smtp_server, smtp_port, sender_email, app_password, recipient_email, subject, body)
    print(review_text)
    return render_template("review.html", username=username)


@app.route('/templates/login.html', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        if authenticate_user(username, password):
            session['username'] = username
            flash('Login successful')
            return redirect('/')
        else:
            flash('Invalid username or password')
    return render_template('login.html')


@app.route('/templates/refresh_password.html', methods=['GET', 'POST'])
def password_recovery():
    return render_template("refresh_password.html")


@app.route('/templates/refresh_password_1.html', methods=['GET', 'POST'])
def password_recovery_2():
    if request.method == 'POST':
        conn = sqlite3.connect("db/users.db")
        cursor = conn.cursor()
        email = request.form.get('email', '').lower().strip()

        smtp_server = "smtp.gmail.com"
        smtp_port = 587
        sender_email = "dictionary.of.evolutions@gmail.com"
        app_password = "jzje ysgq plxm oulg"
        recipient_email = email
        print(recipient_email)

        subject = "Код подтверждения"
        body = str(random.randint(10000, 99999))
        print(body)

        cursor.execute("""
                UPDATE users
                SET code = ?
                WHERE email = ?
            """, (str(body), email))
        conn.commit()
        conn.close()

        send_email(smtp_server, smtp_port, sender_email, app_password, recipient_email, subject, body)
        return render_template("refresh_password_1.html", email=email)
    return render_template("refresh_password.html")


@app.route('/templates/refresh_password_2.html', methods=['GET', 'POST'])
def password_recovery_3():
    if request.method == 'POST':
        conn = sqlite3.connect("db/users.db")
        cursor = conn.cursor()
        code = request.form.get('code', '').strip()
        email = request.form.get('email', '').lower().strip()

        print(f"Email: {email}, Code: {code}")

        cursor.execute("SELECT code FROM users WHERE email = ?", (email,))
        row = cursor.fetchone()
        print(f"Stored code: {row}")

        conn.commit()
        conn.close()
        if code == row[0]:
            return render_template("refresh_password_2.html", email=email)
    return render_template("refresh_password_1.html")


@app.route('/templates/refresh_password_3.html', methods=['GET', 'POST'])
def password_recovery_4():
    if request.method == 'POST':
        password = request.form.get('password', '').strip()
        email = request.form.get('email', '').lower().strip()
        change_password(email, password)
        return render_template("refresh_password_3.html")


@app.route('/templates/sign-up', methods=['GET', 'POST'])
def sign_up():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        email = request.form['email']
        if register_user(username, password, email):
            session['username'] = username
            flash('Registration successful')
            return redirect('/')
        else:
            flash('Username already taken')
    return render_template('sign-up.html')


@app.route('/logout')
def logout():
    session.pop('username', None)
    return redirect('/')


def get_favorits(username):
    conn = sqlite3.connect('db/users.db')
    cursor = conn.cursor()
    cursor.execute("SELECT favorites FROM users WHERE user_name = ?", (username,))
    result = cursor.fetchone()
    conn.close()
    return result[0]


@app.route('/templates/favourites.html')
def favourites():
    username = session.get('username')
    favourit = get_favorits(username)
    if "," in favourit:
        favourit = favourit.split(",")
    else:
        favourit = [favourit]
    print(favourit)
    if favourit == ['']:
        favourit = []
    return render_template('favourites.html', username=username, favourit=favourit)


@app.route('/templates/pictures.html', methods=['GET'])
def pictures_page():
    username = session.get('username')
    image_bool = 0
    item_name = request.args.get('item_name', '').lower().replace(' ', '')
    print(session.get('language'))
    if session:
        update_history(username, item_name)
    conn = sqlite3.connect('db/items.db')
    cursor = conn.cursor()
    if session.get('language') == 'ru':
        cursor.execute("SELECT id, item_name, image, century, custom_text_ru FROM items WHERE item_name LIKE ?",
                       ('%' + item_name + '%',))
    else:
        cursor.execute("SELECT id, item_name, image, century, custom_text FROM items WHERE item_name_en LIKE ?",
                       ('%' + item_name + '%',))
    rows = cursor.fetchall()
    conn.close()
    images = []
    for row in rows:
        image_data = row[2]
        century = row[3]
        image_id = row[0]
        text = row[4]
        if image_data:
            print(f"Image found for item {row[1]}, ID: {image_id}, size: {len(image_data)} bytes, century: {century}")
        else:
            print(f"No image found for item {row[1]} with ID: {image_id}")
        encoded_image_data = base64.b64encode(image_data).decode('utf-8')
        images.append({
            'id': image_id,
            'image_data': encoded_image_data,
            "century": century,
            "text": text
        })
    if images:
        image_bool = 1

    if not images:
        conn = sqlite3.connect('db/images.db')
        cursor = conn.cursor()
        if session.get('language') == 'ru':
            cursor.execute("SELECT id, item_name, image, century, custom_text_ru FROM items WHERE item_name LIKE ?",
                           ('%' + item_name + '%',))
        else:
            cursor.execute("SELECT id, item_name, image, century, custom_text FROM items WHERE item_name_en LIKE ?",
                           ('%' + item_name + '%',))
        rows = cursor.fetchall()

        conn.close()
        images = []
        last_century = 1000
        for row in rows:
            image_data = row[2]
            image_id = row[0]
            century = row[3]
            text = row[4]
            if century == last_century:
                continue
            if image_data:
                print(f"Image found for item {row[1]}, ID: {image_id}, size: {len(image_data)} bytes")
            else:
                print(f"No image found for item {row[1]} with ID: {image_id}")
            encoded_image_data = base64.b64encode(image_data).decode('utf-8')
            images.append({
                'id': image_id,
                'image_data': encoded_image_data,
                "century": century,
                "text": text
            })
            last_century = century
    if not images:
        conn = sqlite3.connect('db/images2.db')
        cursor = conn.cursor()
        if session.get('language') == 'ru':
            cursor.execute("SELECT id, item_name, image, century, custom_text_ru FROM items WHERE item_name LIKE ?",
                           ('%' + item_name + '%',))
        else:
            cursor.execute("SELECT id, item_name, image, century, custom_text FROM items WHERE item_name_en LIKE ?",
                           ('%' + item_name + '%',))
        rows = cursor.fetchall()
        conn.close()
        images = []
        last_century = 1000
        for row in rows:
            image_data = row[2]
            image_id = row[0]
            century = row[3]
            text = row[4]
            if century == last_century:
                continue
            if image_data:
                print(f"Image found for item {row[1]}, ID: {image_id}, size: {len(image_data)} bytes")
            else:
                print(f"No image found for item {row[1]} with ID: {image_id}")
            encoded_image_data = base64.b64encode(image_data).decode('utf-8')
            images.append({
                'id': image_id,
                'image_data': encoded_image_data,
                "century": century,
                "text": text
            })
            last_century = century
    if not images:
        conn = sqlite3.connect('db/items_database (10).db')
        cursor = conn.cursor()
        if session.get('language') == 'ru':
            cursor.execute("SELECT id, item_name, image, century, custom_text_ru FROM items WHERE item_name LIKE ?",
                           ('%' + item_name + '%',))
        else:
            cursor.execute("SELECT id, item_name, image, century, custom_text FROM items WHERE item_name_en LIKE ?",
                           ('%' + item_name + '%',))

        rows = cursor.fetchall()
        conn.close()
        images = []
        last_century = 1000
        for row in rows:
            image_data = row[2]
            image_id = row[0]
            century = row[3]
            text = row[4]
            if century == last_century:
                continue
            if image_data:
                print(f"Image found for item {row[1]}, ID: {image_id}, size: {len(image_data)} bytes")
            else:
                print(f"No image found for item {row[1]} with ID: {image_id}")
            encoded_image_data = base64.b64encode(image_data).decode('utf-8')
            images.append({
                'id': image_id,
                'image_data': encoded_image_data,
                "century": century,
                "text": text
            })
            last_century = century
    if not images:
        conn = sqlite3.connect('db/items_database (11).db')
        cursor = conn.cursor()
        if session.get('language') == 'ru':
            cursor.execute("SELECT id, item_name, image, century, custom_text_ru FROM items WHERE item_name LIKE ?",
                           ('%' + item_name + '%',))
        else:
            cursor.execute("SELECT id, item_name, image, century, custom_text FROM items WHERE item_name_en LIKE ?",
                           ('%' + item_name + '%',))

        rows = cursor.fetchall()
        conn.close()
        images = []
        last_century = 1000
        for row in rows:
            image_data = row[2]
            image_id = row[0]
            century = row[3]
            text = row[4]
            if century == last_century:
                continue
            if image_data:
                print(f"Image found for item {row[1]}, ID: {image_id}, size: {len(image_data)} bytes")
            else:
                print(f"No image found for item {row[1]} with ID: {image_id}")
            encoded_image_data = base64.b64encode(image_data).decode('utf-8')
            images.append({
                'id': image_id,
                'image_data': encoded_image_data,
                "century": century,
                "text": text
            })
            last_century = century
    if images and image_bool == 0:
        image_bool = 2

        # Если пользователь не авторизован, показать страницу без избранного
    if not username:
        if session.get('language') == 'ru':
            return render_template('pictures.html', item_name=translator.translate(item_name.title()), username=username, images=images,
                                   item_text="", image_bool=image_bool,
                                   create_presentation_url=f"/templates/create_of_presentation.html?item_name={item_name}")
        else:
            return render_template('pictures.html', item_name=item_name.title(),
                                   username=username, images=images,
                                   item_text="", image_bool=image_bool,
                                   create_presentation_url=f"/templates/create_of_presentation.html?item_name={item_name}")
    try:
        conn = sqlite3.connect('db/users.db')
        cursor = conn.cursor()
        cursor.execute("SELECT favorites FROM users WHERE user_name = ?", (username,))
        row = cursor.fetchone()
        is_favorite = False
        if row:
            favorites = row[0].split(',') if row[0] else []
            is_favorite = item_name in [i.lower() for i in favorites]
            print(item_name)
            print(favorites)
        print(is_favorite)
        conn.close()

        if session.get('language') == 'ru':
            print(123)
            return render_template('pictures.html', item_name=item_name.title(), username=username, images=images,
                                   item_text="", image_bool=image_bool,
                                   create_presentation_url=f"/templates/create_of_presentation.html?item_name={item_name}",
                                   is_favorite=is_favorite)
        else:
            return render_template('pictures.html', item_name=item_name.title(),
                                   username=username, images=images,
                                   item_text="", image_bool=image_bool,
                                   create_presentation_url=f"/templates/create_of_presentation.html?item_name={item_name}",
                                   is_favorite=is_favorite)
    except Exception as e:
        print(f"Error loading favorites: {e}")
        if session.get('language') == 'ru':
            return render_template('pictures.html', item_name=item_name.title(), username=username, images=images,
                                   item_text="", image_bool=image_bool,
                                   create_presentation_url=f"/templates/create_of_presentation.html?item_name={item_name}",
                                   is_favorite=False)
        else:
            return render_template('pictures.html', item_name=item_name.title(),
                                   username=username, images=images,
                                   item_text="", image_bool=image_bool,
                                   create_presentation_url=f"/templates/create_of_presentation.html?item_name={item_name}",
                                   is_favorite=False)


@app.before_request
def set_default_language():
    if 'language' not in session:
        session['language'] = 'en'


@app.route('/set_language/<lang>')
def set_language(lang):
    if lang in ['en', 'ru']:
        session['language'] = lang
    return redirect(request.referrer or '/')


if __name__ == '__main__':
    initialize_db()
    app.run(debug=True)
